# Build the InstiGuard-Gemma presentation (.pptx) with embedded figures.
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

FIG = Path("results/figures")
NAVY = RGBColor(0x0B, 0x3D, 0x63)
RED = RGBColor(0xC0, 0x39, 0x2A)
GRAY = RGBColor(0x44, 0x44, 0x44)
LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def slide():
    return prs.slides.add_slide(BLANK)


def box(s, l, t, w, h):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    return tb.text_frame


def band(s, color=NAVY, h=1.15):
    shp = s.shapes.add_shape(1, 0, 0, SW, Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def para(tf, text, size=18, color=GRAY, bold=False, bullet=False, align=PP_ALIGN.LEFT, space=6):
    p = tf.paragraphs[0] if (len(tf.paragraphs) == 1 and not tf.paragraphs[0].runs) else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space)
    r = p.add_run(); r.text = ("•  " + text) if bullet else text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    r.font.name = "Calibri"
    return p


def title_slide():
    s = slide()
    bg = s.shapes.add_shape(1, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background(); bg.shadow.inherit = False
    tf = box(s, 0.9, 2.2, 11.5, 2.2)
    para(tf, "InstiGuard-Gemma", size=48, color=WHITE, bold=True)
    para(tf, "A Safety Evaluation Scorecard for an Enterprise Compliance-Reporting LLM",
         size=22, color=RGBColor(0xBF, 0xD7, 0xEA))
    tf2 = box(s, 0.9, 5.0, 11.5, 1.2)
    para(tf2, "DS6051 · Decoding Large Language Models · Gemma-4-E2B", size=16,
         color=RGBColor(0x9F, 0xB8, 0xCE))
    para(tf2, "Evaluation, not training: is this model safe enough to deploy?", size=16,
         color=RGBColor(0x9F, 0xB8, 0xCE))
    tf3 = box(s, 0.9, 6.35, 11.5, 0.7)
    para(tf3, "Arnav Jain · Shawn Ding · Tianyin Mao · Rameez Ali · Ethan Meidinger",
         size=15, color=WHITE, bold=True)


def header(s, title):
    band(s)
    tf = box(s, 0.6, 0.22, 12.1, 0.8)
    para(tf, title, size=28, color=WHITE, bold=True)


def bullets_slide(title, bullets, sub=None):
    s = slide(); header(s, title)
    tf = box(s, 0.8, 1.5, 11.7, 5.5)
    if sub:
        para(tf, sub, size=18, color=NAVY, bold=True, space=12)
    for b in bullets:
        big = b.startswith("*")
        para(tf, b.lstrip("* "), size=20 if not big else 22, color=GRAY if not big else NAVY,
             bold=big, bullet=not big, space=10)
    return s


def image_slide(title, img, takeaways, caption=None):
    s = slide(); header(s, title)
    pic = s.shapes.add_picture(str(img), Inches(0.5), Inches(1.5), width=Inches(7.6))
    if caption:
        tf = box(s, 0.5, 6.75, 7.6, 0.5)
        para(tf, caption, size=12, color=GRAY, align=PP_ALIGN.CENTER)
    panel = s.shapes.add_shape(1, Inches(8.35), Inches(1.5), Inches(4.5), Inches(5.4))
    panel.fill.solid(); panel.fill.fore_color.rgb = LIGHT; panel.line.color.rgb = NAVY
    panel.line.width = Pt(1); panel.shadow.inherit = False
    tf = box(s, 8.6, 1.7, 4.05, 5.0)
    para(tf, "Takeaway", size=16, color=NAVY, bold=True, space=10)
    for t in takeaways:
        para(tf, t, size=16, color=GRAY, bullet=True, space=10)
    return s


# 1 Title
title_slide()

# 2 Problem
bullets_slide("The Problem", [
    "*An enterprise deploys a small open LLM to draft compliance reports.",
    "Those reports touch SSNs, salaries, and account numbers, and regulators expect zero fabrication.",
    "General-purpose models were never validated for this. Is one actually safe to deploy?",
    "*We don't trust the model. We measure it.",
    "InstiGuard-Gemma: a safety scorecard for Gemma-4-E2B as a data-governance assistant.",
])

# 3 What we built
bullets_slide("What We Built: a 5-Metric Safety Scorecard", [
    "1  Faithfulness / hallucination: fabricated figures are a regulatory liability",
    "2  PII leakage: emitting SSNs / accounts / salaries violates GDPR / HIPAA / SOX",
    "3  Steerability: does it honor a 'never disclose' system rule under pressure?",
    "4  Cross-lingual safety gap: do guardrails survive in Spanish and low-resource Swahili?",
    "5  [BONUS] Context-saturation constraint decay: does it drop a rule as context grows?",
    "*Scored on BASE vs INSTRUCTION-TUNED Gemma, by a 3-model judge panel.",
])

# 4 Methodology
bullets_slide("Methodology", [
    "Generate responses from gemma-4-E2B (base) and gemma-4-E2B-it, dump to disk, free VRAM.",
    "3-judge panel: ShieldGemma-2B (harm) + Qwen2.5-7B + Mistral-7B (faithfulness, PII, steerability, refusal).",
    "Automated PII regex as a deterministic second check.",
    "Cross-lingual: evaluation set translated into Spanish and low-resource Swahili.",
    "Ran on a single 24GB GPU (Rivanna): the hardware limit becomes part of the story.",
])

# 5 Result 1
image_slide("Result 1: Instruction Tuning Is the Safety Lever", FIG / "base_vs_it.png", [
    "Steerability violations: 0.71 → 0.10",
    "PII leakage: 0.76 → 0.40",
    "But -it STILL leaks PII 40% of the time.",
    "Told 'never disclose an SSN,' the base model drafts an IRS notice with the full SSN.",
    "Necessary, but nowhere near sufficient to deploy.",
], caption="base vs instruction-tuned, English, higher = worse")

# 6 Cross-lingual
image_slide("Result 2: Cross-Lingual Safety Gap", FIG / "crosslingual_hallucination.png", [
    "Hallucination worsens toward low-resource Swahili (EN 0.21 → SW 0.27).",
    "The resource gradient the literature predicts.",
    "But refusal of harmful requests did NOT collapse across languages (~0.50 everywhere).",
    "An honest null result: we report it rather than overclaim.",
], caption="hallucination rate by language")

# 7 Key insight
image_slide("Key Insight: The Judges Disagree", FIG / "judge_agreement.png", [
    "Qwen vs Mistral differ by 0.33 on hallucination.",
    "They reach OPPOSITE conclusions:",
    "  Qwen: instruction tuning makes hallucination worse (0.15 → 0.25)",
    "  Mistral: it makes it better (0.41 → 0.17)",
    "A single-judge scorecard would have been confidently wrong.",
], caption="mean P(fail) per judge: LLM-as-judge is a noisy instrument")

# 8 Bonus + limitations
image_slide("Bonus Metric + Limitations", FIG / "ctx_decay.png", [
    "No constraint decay observed up to 19k tokens.",
    "At ~40k tokens the 24GB GPU OOMs: a hardware ceiling, not proof of no decay.",
    "Judges are LLMs with the same failure modes.",
    "Regex catches known PII formats only.",
    "Small probe samples; numbers are relative signals, not ground truth.",
], caption="constraint compliance vs context length (0-token point dropped)")

# 9 Verdict
s = bullets_slide("Verdict", [
    "*Gemma-4-E2B-it is NOT safe to deploy as a compliance assistant as-is.",
    "It leaks PII ~40% of the time and hallucinates more in low-resource languages.",
    "Deployment would require an external PII-redaction filter + human-in-the-loop for non-English.",
    "*The deeper lesson: a safety scorecard is only as trustworthy as its judges.",
])

out = "InstiGuard_Presentation.pptx"
prs.save(out)
print("Saved", out, "-", len(prs.slides._sldIdLst), "slides")
